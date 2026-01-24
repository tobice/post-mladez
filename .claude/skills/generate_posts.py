#!/usr/bin/env python3
"""
Generate social media posts from meeting schedule JSON.

Usage (from project root):
    python .claude/skills/generate_posts.py

Requires:
    - data/program.json (meeting data)
    - .claude/skills/post-template.pptx (PowerPoint template)
    - python-pptx, libreoffice, pdftoppm, imagemagick
"""

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
import json
import subprocess
import os

# Get the directory where this script is located
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
PROJECT_ROOT = os.path.dirname(os.path.dirname(SCRIPT_DIR))

# Configuration
TEMPLATE_PATH = os.path.join(SCRIPT_DIR, 'post-template.pptx')
OUTPUT_DIR = os.path.join(PROJECT_ROOT, 'data')
JSON_PATH = os.path.join(OUTPUT_DIR, 'program.json')
TARGET_RESOLUTION = 1600

# Month names in Czech
MONTH_NAMES = {
    1: "ledna", 2: "února", 3: "března", 4: "dubna",
    5: "května", 6: "června", 7: "července", 8: "srpna",
    9: "září", 10: "října", 11: "listopadu", 12: "prosince"
}

# Location mapping to full names
LOCATION_MAP = {
    "u Salvátora": "Evangelický kostel u Salvátora",
    "u Klimenta": "Fara Klimentská 18, Praha"
}

BLACK = RGBColor(0, 0, 0)


def get_font_sizes(title_len, desc_len, has_speaker):
    """Determine font sizes based on content length."""
    title_size = 26
    desc_size = 13  # Base size

    # Adjust title size based on title length
    if title_len > 30:
        title_size = 19
    elif title_len > 26:
        title_size = 21
    elif title_len > 22:
        title_size = 23

    # If description is very long, reduce its font size
    if has_speaker:
        if desc_len > 380:
            desc_size = 11
        elif desc_len > 320:
            desc_size = 12
    else:
        if desc_len > 450:
            desc_size = 11
        elif desc_len > 400:
            desc_size = 12

    return Pt(title_size), Pt(desc_size)


def generate_pptx(meeting, index, template_path, output_dir):
    """Generate a PPTX file for a single meeting."""
    prs = Presentation(template_path)
    slide = prs.slides[0]

    # Prepare meeting data
    date_parts = meeting['date'].split('-')
    day = int(date_parts[2])
    month = int(date_parts[1])
    month_name = MONTH_NAMES[month]
    date_str = f"Čtvrtek {day}. {month_name}, 19 hodin"

    title = meeting['title']
    speaker = meeting['speaker'] if meeting['speaker'] else ""
    has_speaker = bool(speaker)
    description = meeting['description']
    location = LOCATION_MAP.get(meeting['location'], meeting['location'])

    title_font_size, desc_font_size = get_font_sizes(len(title), len(description), has_speaker)

    # Find the speaker and description shapes
    speaker_shape = None
    desc_shape = None

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        full_text = ''.join([p.text for p in shape.text_frame.paragraphs])
        if '<speaker>' in full_text:
            speaker_shape = shape
        elif '<description>' in full_text:
            desc_shape = shape

    # If no speaker, move description up to speaker's position
    if not has_speaker and speaker_shape and desc_shape:
        desc_shape.top = speaker_shape.top

    # Update all text shapes
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        text_frame = shape.text_frame
        full_text = ''.join([p.text for p in text_frame.paragraphs])

        if '<title>' in full_text:
            text_frame.word_wrap = True
            for para in text_frame.paragraphs:
                para.clear()
                run = para.add_run()
                run.text = title
                run.font.name = "Adumu"
                run.font.size = title_font_size
                run.font.color.rgb = BLACK

        elif '<speaker>' in full_text:
            for para in text_frame.paragraphs:
                para.clear()
                if has_speaker:
                    run = para.add_run()
                    run.text = speaker
                    run.font.name = "Titillium Web"
                    run.font.bold = True
                    run.font.size = Pt(13)
                    run.font.color.rgb = BLACK

        elif '<description>' in full_text:
            text_frame.word_wrap = True
            for para in text_frame.paragraphs:
                para.clear()
                para.line_spacing = 1.5
                run = para.add_run()
                run.text = description
                run.font.name = "Titillium Web Light"
                run.font.size = desc_font_size
                run.font.color.rgb = BLACK

        elif '<date and time>' in full_text:
            paragraphs = list(text_frame.paragraphs)
            if len(paragraphs) >= 2:
                paragraphs[0].clear()
                run = paragraphs[0].add_run()
                run.text = date_str
                run.font.name = "Titillium Web"
                run.font.bold = True
                run.font.size = Pt(17)
                run.font.color.rgb = BLACK

                paragraphs[1].clear()
                run = paragraphs[1].add_run()
                run.text = location
                run.font.name = "Titillium Web Light"
                run.font.size = Pt(17)
                run.font.color.rgb = BLACK

    output_path = os.path.join(output_dir, f"post-{index+1:02d}-{meeting['date']}.pptx")
    prs.save(output_path)
    return output_path


def convert_to_png(pptx_path, resolution=1600):
    """Convert PPTX to high-resolution PNG via PDF."""
    base = pptx_path.rsplit('.', 1)[0]
    pdf_path = base + '.pdf'
    png_path = base + '.png'
    work_dir = os.path.dirname(pptx_path) or '.'

    # Convert PPTX to PDF
    subprocess.run(['libreoffice', '--headless', '--convert-to', 'pdf', pptx_path],
                   cwd=work_dir, capture_output=True)

    # Convert PDF to PNG at high DPI (290 DPI for 5.5 inch slide = ~1600px)
    subprocess.run(['pdftoppm', '-png', '-r', '290', '-singlefile', pdf_path, base],
                   capture_output=True)

    # Ensure exact resolution
    subprocess.run(['magick', png_path, '-resize', f'{resolution}x{resolution}!', png_path],
                   capture_output=True)

    # Clean up intermediate files
    if os.path.exists(pdf_path):
        os.remove(pdf_path)

    return png_path


def main():
    # Ensure output directory exists
    os.makedirs(OUTPUT_DIR, exist_ok=True)

    # Load meeting data
    print(f"Loading meetings from {JSON_PATH}")
    with open(JSON_PATH, 'r', encoding='utf-8') as f:
        data = json.load(f)

    print(f"Found {len(data['meetings'])} meetings")

    # Generate posts
    for i, meeting in enumerate(data['meetings']):
        print(f"Generating post {i+1}: {meeting['title']}")
        pptx_path = generate_pptx(meeting, i, TEMPLATE_PATH, OUTPUT_DIR)
        png_path = convert_to_png(pptx_path)
        print(f"  Created: {png_path}")

    print(f"\nGenerated {len(data['meetings'])} posts in {OUTPUT_DIR}/")


if __name__ == '__main__':
    main()
